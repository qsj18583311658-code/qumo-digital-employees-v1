#!/usr/bin/env python3
"""Build a simple editable PowerPoint deck from a JSON spec."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


THEMES = {
    "qumo-clean": {
        "bg": "F8F5F1",
        "ink": "24201E",
        "muted": "756C65",
        "accent": "C55A6A",
        "accent2": "2D7C72",
        "panel": "FFFFFF",
    },
    "editorial-dark": {
        "bg": "171717",
        "ink": "F8F2EA",
        "muted": "C7BEB4",
        "accent": "E0B05F",
        "accent2": "A9D8CE",
        "panel": "242424",
    },
}


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[:2], 16), int(hex_value[2:4], 16), int(hex_value[4:], 16))


def add_textbox(slide, x, y, w, h, text, size, color, bold=False, align=None):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(text)
    if align:
        p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "PingFang SC"
    return shape


def add_bullets(slide, x, y, w, h, bullets, size, color):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "PingFang SC"
        p.space_after = Pt(8)
    return shape


def set_background(slide, theme):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(theme["bg"])


def add_header(slide, theme, title):
    add_textbox(slide, Inches(0.55), Inches(0.32), Inches(11.4), Inches(0.35), "趣摩数字员工", 10, rgb(theme["muted"]))
    add_textbox(slide, Inches(0.55), Inches(0.62), Inches(11.4), Inches(0.8), title, 24, rgb(theme["ink"]), bold=True)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.42), Inches(11.4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(theme["accent"])
    line.line.fill.background()


def add_notes_box(slide, theme, notes):
    if not notes:
        return
    x, y, w, h = Inches(8.5), Inches(6.35), Inches(3.35), Inches(0.65)
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(theme["panel"])
    box.line.color.rgb = rgb(theme["muted"])
    add_textbox(slide, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), h - Inches(0.12), f"讲稿提示：{notes}", 8, rgb(theme["muted"]))


def slide_cover(prs, theme, slide_spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, theme)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(theme["accent"])
    bar.line.fill.background()
    add_textbox(slide, Inches(0.85), Inches(2.0), Inches(10.6), Inches(1.25), slide_spec.get("title", "汇报标题"), 42, rgb(theme["ink"]), True)
    add_textbox(slide, Inches(0.9), Inches(3.35), Inches(9.6), Inches(0.7), slide_spec.get("subtitle", ""), 18, rgb(theme["muted"]))
    add_textbox(slide, Inches(0.9), Inches(6.55), Inches(9.6), Inches(0.35), slide_spec.get("footer", "趣摩文化"), 11, rgb(theme["muted"]))
    return slide


def slide_summary(prs, theme, slide_spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, theme)
    add_header(slide, theme, slide_spec.get("title", "核心结论"))
    bullets = slide_spec.get("bullets") or slide_spec.get("points") or []
    add_bullets(slide, Inches(0.85), Inches(1.9), Inches(11.0), Inches(4.8), bullets, 22, rgb(theme["ink"]))
    add_notes_box(slide, theme, slide_spec.get("speaker_note"))
    return slide


def slide_metrics(prs, theme, slide_spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, theme)
    add_header(slide, theme, slide_spec.get("title", "关键指标"))
    metrics = slide_spec.get("metrics") or []
    if not metrics:
        metrics = [{"label": "待补数据", "value": "待补数据", "note": "请补充平台导出或活动表"}]
    cols = min(4, len(metrics))
    card_w = 10.9 / cols
    for i, metric in enumerate(metrics[:8]):
        row = i // cols
        col = i % cols
        x = Inches(0.75 + col * card_w)
        y = Inches(1.95 + row * 2.05)
        box = slide.shapes.add_shape(1, x, y, Inches(card_w - 0.18), Inches(1.55))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(theme["panel"])
        box.line.color.rgb = rgb(theme["accent"])
        add_textbox(slide, x + Inches(0.16), y + Inches(0.12), Inches(card_w - 0.5), Inches(0.32), metric.get("label", ""), 12, rgb(theme["muted"]))
        add_textbox(slide, x + Inches(0.16), y + Inches(0.48), Inches(card_w - 0.5), Inches(0.48), metric.get("value", ""), 25, rgb(theme["ink"]), True)
        add_textbox(slide, x + Inches(0.16), y + Inches(1.05), Inches(card_w - 0.5), Inches(0.3), metric.get("note", ""), 9, rgb(theme["muted"]))
    add_notes_box(slide, theme, slide_spec.get("speaker_note"))
    return slide


def slide_two_column(prs, theme, slide_spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, theme)
    add_header(slide, theme, slide_spec.get("title", "对比分析"))
    add_textbox(slide, Inches(0.85), Inches(1.85), Inches(5.2), Inches(0.45), slide_spec.get("left_title", "左侧判断"), 18, rgb(theme["accent"]), True)
    add_bullets(slide, Inches(0.85), Inches(2.35), Inches(5.1), Inches(4.0), slide_spec.get("left") or [], 16, rgb(theme["ink"]))
    add_textbox(slide, Inches(6.55), Inches(1.85), Inches(5.2), Inches(0.45), slide_spec.get("right_title", "右侧判断"), 18, rgb(theme["accent2"]), True)
    add_bullets(slide, Inches(6.55), Inches(2.35), Inches(5.1), Inches(4.0), slide_spec.get("right") or [], 16, rgb(theme["ink"]))
    add_notes_box(slide, theme, slide_spec.get("speaker_note"))
    return slide


def slide_actions(prs, theme, slide_spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, theme)
    add_header(slide, theme, slide_spec.get("title", "下一步行动"))
    actions = slide_spec.get("actions") or []
    y = Inches(1.8)
    headers = ["动作", "负责人", "时间", "指标"]
    widths = [5.2, 1.8, 1.7, 2.4]
    x0 = Inches(0.75)
    x = x0
    for header, width in zip(headers, widths):
        add_textbox(slide, x, y, Inches(width), Inches(0.35), header, 11, rgb(theme["accent"]), True)
        x += Inches(width)
    y += Inches(0.45)
    for action in actions[:8]:
        x = x0
        vals = [
            action.get("item", ""),
            action.get("owner", ""),
            action.get("deadline", ""),
            action.get("metric", ""),
        ]
        for val, width in zip(vals, widths):
            add_textbox(slide, x, y, Inches(width - 0.1), Inches(0.48), val, 11, rgb(theme["ink"]))
            x += Inches(width)
        y += Inches(0.62)
    add_notes_box(slide, theme, slide_spec.get("speaker_note"))
    return slide


def normalize_slides(spec: dict[str, Any]) -> list[dict[str, Any]]:
    slides = spec.get("slides") or []
    if slides:
        return slides
    return [
        {"type": "cover", "title": spec.get("title", "汇报标题"), "subtitle": spec.get("subtitle", "")},
        {"type": "summary", "title": "核心结论", "bullets": spec.get("bullets") or ["待补结论 1", "待补结论 2", "待补结论 3"]},
        {"type": "actions", "title": "下一步行动", "actions": spec.get("actions") or []},
    ]


def main() -> None:
    parser = argparse.ArgumentParser(description="Build an editable PPTX deck from JSON.")
    parser.add_argument("spec", type=Path, help="JSON spec path")
    parser.add_argument("--out", type=Path, default=Path("qumo_deck.pptx"), help="output PPTX path")
    args = parser.parse_args()

    spec = json.loads(args.spec.read_text(encoding="utf-8"))
    theme = THEMES.get(spec.get("theme", "qumo-clean"), THEMES["qumo-clean"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = {
        "cover": slide_cover,
        "summary": slide_summary,
        "metrics": slide_metrics,
        "two-column": slide_two_column,
        "actions": slide_actions,
    }
    for slide_spec in normalize_slides(spec):
        builder = builders.get(slide_spec.get("type", "summary"), slide_summary)
        builder(prs, theme, slide_spec)

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(json.dumps({"generated": str(args.out), "slides": len(prs.slides)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
